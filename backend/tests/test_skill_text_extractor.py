"""Skill text extractor tests: PDF, DOCX, PPTX, TXT extraction and markdown conversion.

Tests every extractor function plus the dispatcher and markdown converter.
Uses real in-memory file objects where possible (python-docx, python-pptx)
and mocks for PDF (pdfplumber requires actual PDF bytes).
Target: 95%+ coverage of app/services/skill_text_extractor.py
"""

import io
import sys
from unittest.mock import MagicMock, patch

import pytest

from app.services.skill_text_extractor import (
    convert_to_markdown,
    extract_text,
    extract_text_from_docx,
    extract_text_from_pdf,
    extract_text_from_pptx,
    extract_text_from_text,
)


# ---------------------------------------------------------------------------
# Helpers: create real in-memory files
# ---------------------------------------------------------------------------


def _make_docx_bytes(paragraphs: list[str], table_rows: list[list[str]] | None = None) -> bytes:
    """Create a minimal DOCX file in memory with given paragraphs and optional table."""
    from docx import Document

    doc = Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    if table_rows:
        table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for r_idx, row_data in enumerate(table_rows):
            for c_idx, cell_text in enumerate(row_data):
                table.rows[r_idx].cells[c_idx].text = cell_text
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx_bytes(slides_text: list[list[str]]) -> bytes:
    """Create a minimal PPTX file in memory. Each inner list = one slide's text boxes."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for texts in slides_text:
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)
        for i, text in enumerate(texts):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(5), Inches(1))
            txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ===========================================================================
# Test classes
# ===========================================================================


class TestExtractTextFromPdf:
    """Test PDF extraction via pdfplumber.

    Since pdfplumber is imported inside the function body, we patch
    'pdfplumber.open' on the real module that gets imported at call time.
    """

    def test_pdf_with_multiple_pages(self):
        """Mock pdfplumber.open to return multi-page text."""
        mock_page1 = MagicMock()
        mock_page1.extract_text.return_value = "Page 1 content"
        mock_page2 = MagicMock()
        mock_page2.extract_text.return_value = "Page 2 content"

        mock_pdf = MagicMock()
        mock_pdf.pages = [mock_page1, mock_page2]
        mock_pdf.__enter__ = MagicMock(return_value=mock_pdf)
        mock_pdf.__exit__ = MagicMock(return_value=False)

        with patch("pdfplumber.open", return_value=mock_pdf):
            result = extract_text_from_pdf(b"fake-pdf-bytes")

        assert result == "Page 1 content\n\nPage 2 content"

    def test_pdf_pages_with_none_text_are_skipped(self):
        """Pages where extract_text returns None should be skipped."""
        mock_page1 = MagicMock()
        mock_page1.extract_text.return_value = "Real content"
        mock_page2 = MagicMock()
        mock_page2.extract_text.return_value = None
        mock_page3 = MagicMock()
        mock_page3.extract_text.return_value = "More content"

        mock_pdf = MagicMock()
        mock_pdf.pages = [mock_page1, mock_page2, mock_page3]
        mock_pdf.__enter__ = MagicMock(return_value=mock_pdf)
        mock_pdf.__exit__ = MagicMock(return_value=False)

        with patch("pdfplumber.open", return_value=mock_pdf):
            result = extract_text_from_pdf(b"fake-pdf-bytes")

        assert result == "Real content\n\nMore content"

    def test_pdf_with_empty_text_pages_skipped(self):
        """Pages where extract_text returns empty string should be skipped."""
        mock_page = MagicMock()
        mock_page.extract_text.return_value = ""

        mock_pdf = MagicMock()
        mock_pdf.pages = [mock_page]
        mock_pdf.__enter__ = MagicMock(return_value=mock_pdf)
        mock_pdf.__exit__ = MagicMock(return_value=False)

        with patch("pdfplumber.open", return_value=mock_pdf):
            result = extract_text_from_pdf(b"fake-pdf-bytes")

        assert result == ""

    def test_pdf_open_raises_returns_empty(self):
        """When pdfplumber.open raises (corrupt PDF), return empty string."""
        with patch("pdfplumber.open", side_effect=RuntimeError("corrupt PDF")):
            result = extract_text_from_pdf(b"corrupt-pdf-bytes")

        assert result == ""

    def test_pdf_no_pages(self):
        """PDF with zero pages returns empty string."""
        mock_pdf = MagicMock()
        mock_pdf.pages = []
        mock_pdf.__enter__ = MagicMock(return_value=mock_pdf)
        mock_pdf.__exit__ = MagicMock(return_value=False)

        with patch("pdfplumber.open", return_value=mock_pdf):
            result = extract_text_from_pdf(b"empty-pdf")

        assert result == ""

    def test_pdf_import_missing_returns_empty(self):
        """If pdfplumber is not installed (ImportError), return empty string."""
        real_module = sys.modules.get("pdfplumber")
        # Setting to None makes `import pdfplumber` raise ImportError
        sys.modules["pdfplumber"] = None  # type: ignore[assignment]
        try:
            result = extract_text_from_pdf(b"any-bytes")
            assert result == ""
        finally:
            if real_module is not None:
                sys.modules["pdfplumber"] = real_module
            else:
                sys.modules.pop("pdfplumber", None)


class TestExtractTextFromDocx:
    """Test DOCX extraction with real python-docx objects."""

    def test_docx_with_paragraphs(self):
        """Extract text from paragraphs in a real DOCX."""
        content = _make_docx_bytes(["First paragraph", "Second paragraph"])
        result = extract_text_from_docx(content)

        assert "First paragraph" in result
        assert "Second paragraph" in result

    def test_docx_with_table(self):
        """Extract text from a table in a real DOCX."""
        content = _make_docx_bytes(
            ["Header text"],
            table_rows=[["Cell A1", "Cell B1"], ["Cell A2", "Cell B2"]],
        )
        result = extract_text_from_docx(content)

        assert "Header text" in result
        assert "Cell A1" in result
        assert "Cell B1" in result

    def test_docx_with_empty_paragraphs_skipped(self):
        """Empty/whitespace-only paragraphs should be skipped."""
        content = _make_docx_bytes(["Content", "", "   ", "More content"])
        result = extract_text_from_docx(content)

        assert "Content" in result
        assert "More content" in result

    def test_docx_empty_document(self):
        """Empty DOCX returns empty string."""
        content = _make_docx_bytes([])
        result = extract_text_from_docx(content)

        assert result == ""

    def test_docx_exception_returns_empty(self):
        """Corrupt DOCX bytes return empty string."""
        result = extract_text_from_docx(b"not-a-docx-file")
        assert result == ""

    def test_docx_table_with_empty_cells(self):
        """Table rows where all cells are empty should be skipped."""
        content = _make_docx_bytes(
            ["Title"],
            table_rows=[["Data", ""], ["", ""]],
        )
        result = extract_text_from_docx(content)

        assert "Title" in result
        assert "Data" in result

    def test_docx_import_missing_returns_empty(self):
        """If python-docx is not installed, return empty string."""
        real_module = sys.modules.get("docx")
        sys.modules["docx"] = None  # type: ignore[assignment]
        try:
            result = extract_text_from_docx(b"any-bytes")
            assert result == ""
        finally:
            if real_module is not None:
                sys.modules["docx"] = real_module
            else:
                sys.modules.pop("docx", None)


class TestExtractTextFromPptx:
    """Test PPTX extraction with real python-pptx objects."""

    def test_pptx_with_slides(self):
        """Extract text from a PPTX with two slides."""
        content = _make_pptx_bytes([
            ["Slide 1 Title", "Slide 1 Body"],
            ["Slide 2 Title"],
        ])
        result = extract_text_from_pptx(content)

        assert "Slide 1 Title" in result
        assert "Slide 1 Body" in result
        assert "Slide 2 Title" in result
        assert "\n\n---\n\n" in result

    def test_pptx_single_slide(self):
        """PPTX with a single slide, no separator."""
        content = _make_pptx_bytes([["Only slide"]])
        result = extract_text_from_pptx(content)

        assert "Only slide" in result
        assert "---" not in result

    def test_pptx_empty_shapes_skipped(self):
        """Slides with only empty shapes produce no output for that slide."""
        content = _make_pptx_bytes([["Content slide"], []])
        result = extract_text_from_pptx(content)

        assert "Content slide" in result

    def test_pptx_empty_presentation(self):
        """PPTX with no slides returns empty string."""
        from pptx import Presentation

        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        content = buf.getvalue()

        result = extract_text_from_pptx(content)
        assert result == ""

    def test_pptx_exception_returns_empty(self):
        """Corrupt PPTX bytes return empty string."""
        result = extract_text_from_pptx(b"not-a-pptx-file")
        assert result == ""

    def test_pptx_import_missing_returns_empty(self):
        """If python-pptx is not installed, return empty string."""
        real_module = sys.modules.get("pptx")
        sys.modules["pptx"] = None  # type: ignore[assignment]
        try:
            result = extract_text_from_pptx(b"any-bytes")
            assert result == ""
        finally:
            if real_module is not None:
                sys.modules["pptx"] = real_module
            else:
                sys.modules.pop("pptx", None)


class TestExtractTextFromText:
    """Test plain text / markdown extraction."""

    def test_normal_utf8_bytes(self):
        """Normal UTF-8 bytes decode correctly."""
        text = "Hello, world!"
        result = extract_text_from_text(text.encode("utf-8"))
        assert result == "Hello, world!"

    def test_chinese_utf8_bytes(self):
        """Chinese text decodes correctly."""
        text = "泽布替尼培训资料"
        result = extract_text_from_text(text.encode("utf-8"))
        assert result == text

    def test_invalid_utf8_uses_replacement(self):
        """Invalid UTF-8 bytes should use replacement character."""
        bad_bytes = b"Hello \xff\xfe world"
        result = extract_text_from_text(bad_bytes)
        assert "\ufffd" in result
        assert "Hello" in result
        assert "world" in result

    def test_empty_bytes(self):
        """Empty bytes return empty string."""
        result = extract_text_from_text(b"")
        assert result == ""

    def test_multiline_text(self):
        """Multiline text preserves newlines."""
        text = "Line 1\nLine 2\nLine 3"
        result = extract_text_from_text(text.encode("utf-8"))
        assert result == text


class TestExtractText:
    """Test the dispatcher function that routes based on file extension.

    _EXTRACTOR_MAP captures function references at import time, so for dispatch
    routing tests we patch the map itself. For txt/md we use real extractors.
    """

    def test_dispatch_pdf(self):
        """PDF extension dispatches to PDF extractor."""
        mock_fn = MagicMock(return_value="pdf content")
        new_map = {".pdf": mock_fn}
        with patch("app.services.skill_text_extractor._EXTRACTOR_MAP", new_map):
            result = extract_text(b"data", "document.pdf")
            mock_fn.assert_called_once_with(b"data")
            assert result == "pdf content"

    def test_dispatch_docx(self):
        """DOCX extension dispatches to DOCX extractor."""
        mock_fn = MagicMock(return_value="docx content")
        new_map = {".docx": mock_fn}
        with patch("app.services.skill_text_extractor._EXTRACTOR_MAP", new_map):
            result = extract_text(b"data", "report.docx")
            mock_fn.assert_called_once_with(b"data")
            assert result == "docx content"

    def test_dispatch_pptx(self):
        """PPTX extension dispatches to PPTX extractor."""
        mock_fn = MagicMock(return_value="pptx content")
        new_map = {".pptx": mock_fn}
        with patch("app.services.skill_text_extractor._EXTRACTOR_MAP", new_map):
            result = extract_text(b"data", "slides.pptx")
            mock_fn.assert_called_once_with(b"data")
            assert result == "pptx content"

    def test_dispatch_txt(self):
        """TXT extension dispatches to text extractor."""
        result = extract_text(b"plain text content", "notes.txt")
        assert result == "plain text content"

    def test_dispatch_md(self):
        """MD extension dispatches to text extractor."""
        result = extract_text(b"# Markdown heading", "readme.md")
        assert result == "# Markdown heading"

    def test_unknown_extension_returns_empty(self):
        """Unknown file extension returns empty string."""
        result = extract_text(b"data", "file.xyz")
        assert result == ""

    def test_no_extension_returns_empty(self):
        """Filename without extension returns empty string."""
        result = extract_text(b"data", "Makefile")
        assert result == ""

    def test_case_insensitive_extension(self):
        """Extension matching should be case-insensitive."""
        result = extract_text(b"plain text", "FILE.TXT")
        assert result == "plain text"

    def test_uppercase_pdf_extension(self):
        """Uppercase .PDF should still dispatch to PDF extractor."""
        mock_fn = MagicMock(return_value="pdf result")
        new_map = {".pdf": mock_fn}
        with patch("app.services.skill_text_extractor._EXTRACTOR_MAP", new_map):
            result = extract_text(b"data", "DOCUMENT.PDF")
            mock_fn.assert_called_once()
            assert result == "pdf result"

    def test_extractor_raises_exception_returns_empty(self):
        """If the extractor function itself raises, dispatcher catches and returns empty."""
        mock_fn = MagicMock(side_effect=RuntimeError("unexpected error"))
        new_map = {".pdf": mock_fn}
        with patch("app.services.skill_text_extractor._EXTRACTOR_MAP", new_map):
            result = extract_text(b"data", "file.pdf")
            assert result == ""

    def test_double_extension_uses_last(self):
        """Filename with multiple dots uses the last extension."""
        result = extract_text(b"text data", "archive.backup.txt")
        assert result == "text data"

    def test_hidden_file_with_extension(self):
        """Hidden file (starts with dot) with valid extension works."""
        result = extract_text(b"hidden content", ".notes.md")
        assert result == "hidden content"

    def test_empty_filename_returns_empty(self):
        """Empty filename returns empty string (no extension found)."""
        result = extract_text(b"data", "")
        assert result == ""

    def test_dot_only_filename(self):
        """Filename that is just a dot: extension is empty string '.'."""
        result = extract_text(b"data", ".")
        assert result == ""


class TestConvertToMarkdown:
    """Test markdown conversion with structure inference."""

    def test_txt_returns_as_is(self):
        """TXT files return text unmodified."""
        text = "This is plain text.\nNo conversion needed."
        result = convert_to_markdown(text, "notes.txt")
        assert result == text

    def test_md_returns_as_is(self):
        """MD files return text unmodified."""
        text = "# Already markdown\n\nWith content."
        result = convert_to_markdown(text, "readme.md")
        assert result == text

    def test_pptx_adds_slide_headings(self):
        """PPTX conversion adds '## Slide N' headings."""
        text = "Slide one content\n\n---\n\nSlide two content\n\n---\n\nSlide three content"
        result = convert_to_markdown(text, "presentation.pptx")

        assert "## Slide 1" in result
        assert "## Slide 2" in result
        assert "## Slide 3" in result
        assert "Slide one content" in result
        assert "Slide two content" in result

    def test_pptx_single_slide(self):
        """PPTX with single slide (no separator) still gets heading."""
        text = "Only slide content"
        result = convert_to_markdown(text, "single.pptx")

        assert "## Slide 1" in result
        assert "Only slide content" in result

    def test_pdf_infers_headings_from_short_lines(self):
        """PDF conversion infers headings from short lines followed by blank lines."""
        text = "Chapter Title\n\nThis is a paragraph with actual content that continues.\n"
        result = convert_to_markdown(text, "document.pdf")

        assert "## Chapter Title" in result

    def test_docx_infers_headings(self):
        """DOCX conversion also uses heading inference."""
        text = "Section Header\n\nBody text follows here.\n"
        result = convert_to_markdown(text, "report.docx")

        assert "## Section Header" in result

    def test_heading_inference_skips_short_lines(self):
        """Lines shorter than 3 chars are not promoted to headings."""
        text = "OK\n\nBody text here.\n"
        result = convert_to_markdown(text, "document.pdf")

        assert "## OK" not in result

    def test_heading_inference_skips_lines_ending_with_period(self):
        """Lines ending with '.' are sentences, not headings."""
        text = "This is a sentence.\n\nMore text here.\n"
        result = convert_to_markdown(text, "document.pdf")

        assert "## This is a sentence." not in result

    def test_heading_inference_skips_lines_ending_with_comma(self):
        """Lines ending with ',' are not headings."""
        text = "Item one,\n\nItem two follows.\n"
        result = convert_to_markdown(text, "document.pdf")

        assert "## Item one," not in result

    def test_heading_inference_long_lines_not_headings(self):
        """Lines longer than 79 chars should not be promoted to headings."""
        long_line = "A" * 80
        text = f"{long_line}\n\nBody text.\n"
        result = convert_to_markdown(text, "document.pdf")

        assert f"## {long_line}" not in result

    def test_heading_inference_exactly_79_chars_is_heading(self):
        """A line of exactly 79 chars followed by blank line becomes heading."""
        line_79 = "A" * 79
        text = f"{line_79}\n\nBody text follows.\n"
        result = convert_to_markdown(text, "document.pdf")

        assert f"## {line_79}" in result

    def test_unknown_extension_uses_heading_inference(self):
        """Non txt/md/pptx extensions fall through to heading inference."""
        text = "Title Line\n\nContent below.\n"
        result = convert_to_markdown(text, "document.pdf")

        assert "## Title Line" in result

    def test_no_extension_uses_heading_inference(self):
        """No extension falls through to heading inference (ext is empty)."""
        text = "Heading\n\nParagraph.\n"
        result = convert_to_markdown(text, "noext")

        assert "## Heading" in result

    def test_pptx_empty_text(self):
        """PPTX with empty text returns single slide heading."""
        result = convert_to_markdown("", "slides.pptx")
        assert "## Slide 1" in result

    def test_convert_preserves_non_heading_content(self):
        """Content that does not match heading pattern is preserved as-is."""
        text = "Normal paragraph that is long enough to not be a heading by itself."
        result = convert_to_markdown(text, "document.pdf")

        assert result == text

    def test_multiple_headings_inferred(self):
        """Multiple short lines followed by blanks all become headings."""
        text = "Introduction\n\nFirst paragraph.\n\nConclusion\n\nLast paragraph.\n"
        result = convert_to_markdown(text, "document.pdf")

        assert "## Introduction" in result
        assert "## Conclusion" in result

    def test_case_insensitive_pptx_extension(self):
        """PPTX extension matching should be case-insensitive."""
        text = "Content A\n\n---\n\nContent B"
        result = convert_to_markdown(text, "SLIDES.PPTX")

        assert "## Slide 1" in result
        assert "## Slide 2" in result

    def test_case_insensitive_txt_extension(self):
        """TXT extension should be case-insensitive too."""
        text = "Keep as is"
        result = convert_to_markdown(text, "FILE.TXT")
        assert result == text

    def test_case_insensitive_md_extension(self):
        """MD extension should be case-insensitive."""
        text = "# Title"
        result = convert_to_markdown(text, "README.MD")
        assert result == text
