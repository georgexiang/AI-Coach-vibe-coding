"""Text extraction service for converting uploaded materials to plain text.

Supports PDF, DOCX, PPTX, TXT, and MD files. Each extractor wraps its
processing in try/except for error resilience -- a single bad file never
crashes the batch.
"""

import io
import logging
import re

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------


def extract_text_from_pdf(content: bytes) -> str:
    """Extract text from a PDF file using pdfplumber (better table support)."""
    try:
        import pdfplumber

        pages_text: list[str] = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text)
        return "\n\n".join(pages_text)
    except Exception as exc:
        logger.warning("PDF extraction failed: %s", exc)
        return ""


def extract_text_from_docx(content: bytes) -> str:
    """Extract text from a DOCX file, including paragraphs and tables."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(content))

        parts: list[str] = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)

        # Extract table content
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        return "\n".join(parts)
    except Exception as exc:
        logger.warning("DOCX extraction failed: %s", exc)
        return ""


def extract_text_from_pptx(content: bytes) -> str:
    """Extract text from a PPTX file, one section per slide."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slides_text: list[str] = []

        for slide in prs.slides:
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text)
            if slide_parts:
                slides_text.append("\n".join(slide_parts))

        return "\n\n---\n\n".join(slides_text)
    except Exception as exc:
        logger.warning("PPTX extraction failed: %s", exc)
        return ""


def extract_text_from_text(content: bytes) -> str:
    """Decode plain text / markdown content (UTF-8 with replacement)."""
    return content.decode("utf-8", errors="replace")


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

_EXTRACTOR_MAP: dict[str, object] = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".pptx": extract_text_from_pptx,
    ".txt": extract_text_from_text,
    ".md": extract_text_from_text,
}


def extract_text(content: bytes, filename: str) -> str:
    """Dispatch to the appropriate extractor based on file extension.

    Returns extracted text, or empty string on failure (never raises).
    """
    ext = ""
    dot_idx = filename.rfind(".")
    if dot_idx >= 0:
        ext = filename[dot_idx:].lower()

    extractor = _EXTRACTOR_MAP.get(ext)
    if extractor is None:
        logger.warning("No extractor for extension '%s' (file: %s)", ext, filename)
        return ""

    try:
        return extractor(content)  # type: ignore[operator]
    except Exception as exc:
        logger.error("Text extraction failed for %s: %s", filename, exc)
        return ""


# ---------------------------------------------------------------------------
# Markdown conversion helpers
# ---------------------------------------------------------------------------

# Regex: lines shorter than 80 chars followed by a blank line -> heading
_HEADING_RE = re.compile(r"^(.{1,79})$\n\n", re.MULTILINE)


def convert_to_markdown(text: str, filename: str) -> str:
    """Convert extracted text to Markdown with lightweight structure inference.

    - PDF/DOCX: infer headings from short lines followed by blank lines.
    - PPTX: label each slide section.
    - TXT/MD: return as-is.
    """
    ext = ""
    dot_idx = filename.rfind(".")
    if dot_idx >= 0:
        ext = filename[dot_idx:].lower()

    if ext in (".txt", ".md"):
        return text

    if ext == ".pptx":
        sections = text.split("\n\n---\n\n")
        md_parts: list[str] = []
        for idx, section in enumerate(sections, 1):
            md_parts.append(f"## Slide {idx}\n\n{section}")
        return "\n\n".join(md_parts)

    # PDF / DOCX: infer headings
    def _heading_replacer(match: re.Match[str]) -> str:
        line = match.group(1).strip()
        # Skip very short lines (likely page numbers) or lines that look like sentences
        if len(line) < 3 or line.endswith(".") or line.endswith(","):
            return match.group(0)
        return f"## {line}\n\n"

    return _HEADING_RE.sub(_heading_replacer, text)
